"""
pptx_builder.py
---------------
PowerPoint utility functions for GDD pitch deck generation using python-pptx.
Provides slide layout templates, text box creation, color theme management,
bullet list rendering, and slide builders for title, content, and comparison slides.
"""

from typing import List, Tuple, Optional, Dict, Any
from dataclasses import dataclass, field

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Emu
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    Presentation = Inches = Pt = Emu = RGBColor = PP_ALIGN = None


# ─────────────────────────────────────────────
# THEME DEFINITION
# ─────────────────────────────────────────────

@dataclass
class PitchTheme:
    """Color and font theme for the pitch deck."""
    # Backgrounds
    bg_dark: Tuple[int, int, int] = (15, 23, 42)         # Dark navy
    bg_slide: Tuple[int, int, int] = (22, 33, 57)         # Slightly lighter navy
    bg_card: Tuple[int, int, int] = (30, 45, 78)          # Card/section bg
    bg_accent: Tuple[int, int, int] = (41, 89, 133)       # Mid blue accent

    # Text
    text_primary: Tuple[int, int, int] = (240, 244, 255)  # Near white
    text_secondary: Tuple[int, int, int] = (160, 185, 215) # Muted blue-white
    text_accent: Tuple[int, int, int] = (232, 184, 75)    # Gold for highlights

    # Accents
    accent_gold: Tuple[int, int, int] = (232, 184, 75)    # Gold
    accent_blue: Tuple[int, int, int] = (66, 153, 225)    # Bright blue
    accent_green: Tuple[int, int, int] = (72, 187, 120)   # Success green
    accent_red: Tuple[int, int, int] = (245, 101, 101)    # Warning red

    # Fonts
    font_heading: str = "Calibri"
    font_body: str = "Calibri"

    # Slide dimensions (16:9 widescreen)
    slide_width: float = 10.0   # inches
    slide_height: float = 5.625  # inches


# Default theme
DEFAULT_THEME = PitchTheme()


# ─────────────────────────────────────────────
# HELPERS
# ─────────────────────────────────────────────

def rgb(color_tuple: Tuple[int, int, int]) -> "RGBColor":
    """Convert (r, g, b) tuple to pptx RGBColor."""
    if not PPTX_AVAILABLE:
        return None
    return RGBColor(color_tuple[0], color_tuple[1], color_tuple[2])


def create_presentation(theme: PitchTheme = DEFAULT_THEME) -> "Presentation":
    """Create a new blank presentation with theme dimensions."""
    if not PPTX_AVAILABLE:
        raise ImportError("python-pptx is required. Install with: pip install python-pptx")
    prs = Presentation()
    prs.slide_width = Inches(theme.slide_width)
    prs.slide_height = Inches(theme.slide_height)
    return prs


def add_background(slide, color: Tuple[int, int, int], theme: PitchTheme = DEFAULT_THEME) -> None:
    """Fill slide background with a solid color."""
    if not PPTX_AVAILABLE:
        return
    from pptx.util import Inches
    from pptx.enum.dml import MSO_THEME_COLOR
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(color)


def add_text_box(
    slide,
    text: str,
    left: float, top: float, width: float, height: float,
    font_name: str = "Calibri",
    font_size: int = 14,
    bold: bool = False,
    italic: bool = False,
    color: Tuple[int, int, int] = (240, 244, 255),
    align: str = "left",
    word_wrap: bool = True
) -> Any:
    """
    Add a styled text box to a slide.

    Args:
        slide: python-pptx slide object
        text: Text content
        left, top, width, height: Position and size in inches
        font_name: Font family name
        font_size: Font size in points
        bold, italic: Text styling
        color: RGB color tuple
        align: "left", "center", "right"
        word_wrap: Enable word wrap

    Returns:
        TextBox shape object
    """
    if not PPTX_AVAILABLE:
        return None
    from pptx.util import Inches, Pt
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = word_wrap

    p = tf.paragraphs[0]
    p.text = text

    if align == "center":
        p.alignment = PP_ALIGN.CENTER
    elif align == "right":
        p.alignment = PP_ALIGN.RIGHT
    else:
        p.alignment = PP_ALIGN.LEFT

    run = p.runs[0] if p.runs else p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = rgb(color)

    return txBox


def add_bullet_text_box(
    slide,
    title: str,
    bullets: List[str],
    left: float, top: float, width: float, height: float,
    theme: PitchTheme = DEFAULT_THEME,
    title_size: int = 16,
    bullet_size: int = 12
) -> Any:
    """
    Add a text box with a title and bulleted list.

    Args:
        slide: Slide object
        title: Box heading
        bullets: List of bullet point strings
        left, top, width, height: Position in inches
        theme: Color theme
        title_size, bullet_size: Font sizes

    Returns:
        TextBox shape
    """
    if not PPTX_AVAILABLE:
        return None
    from pptx.util import Inches, Pt
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    # Title paragraph
    p = tf.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0] if p.runs else p.add_run()
    run.text = title
    run.font.name = theme.font_heading
    run.font.size = Pt(title_size)
    run.font.bold = True
    run.font.color.rgb = rgb(theme.accent_gold)

    # Bullet paragraphs
    for bullet in bullets:
        p_new = tf.add_paragraph()
        p_new.text = f"  •  {bullet}"
        p_new.alignment = PP_ALIGN.LEFT
        run_b = p_new.runs[0] if p_new.runs else p_new.add_run()
        run_b.text = f"  •  {bullet}"
        run_b.font.name = theme.font_body
        run_b.font.size = Pt(bullet_size)
        run_b.font.color.rgb = rgb(theme.text_primary)

    return txBox


def add_accent_bar(
    slide,
    left: float, top: float, width: float, height: float,
    color: Tuple[int, int, int]
) -> None:
    """Add a solid color rectangle (used for accent bars and section markers)."""
    if not PPTX_AVAILABLE:
        return
    from pptx.util import Inches
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    shape.line.fill.background()  # No border


def add_divider_line(
    slide,
    x1: float, y1: float, x2: float, y2: float,
    color: Tuple[int, int, int] = (41, 89, 133),
    width_pt: float = 1.0
) -> None:
    """Add a horizontal or vertical line to a slide."""
    if not PPTX_AVAILABLE:
        return
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    line = slide.shapes.add_connector(
        1,  # STRAIGHT connector
        Inches(x1), Inches(y1),
        Inches(x2), Inches(y2)
    )
    line.line.color.rgb = rgb(color)
    line.line.width = Pt(width_pt)


# ─────────────────────────────────────────────
# SLIDE BUILDERS
# ─────────────────────────────────────────────

def build_title_slide(
    prs: "Presentation",
    game_title: str,
    tagline: str,
    studio_name: str,
    genre: str,
    platform: str,
    theme: PitchTheme = DEFAULT_THEME
) -> None:
    """
    Build the title slide (Slide 1).
    Layout: Full dark background, large title, tagline, studio name, metadata.
    """
    if not PPTX_AVAILABLE:
        return
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    add_background(slide, theme.bg_dark)

    # Left accent bar
    add_accent_bar(slide, 0, 0, 0.08, theme.slide_height, theme.accent_gold)

    # Game title
    add_text_box(
        slide, game_title,
        0.3, 1.2, 9.5, 1.8,
        font_name=theme.font_heading,
        font_size=48, bold=True,
        color=theme.text_primary, align="left"
    )

    # Tagline
    add_text_box(
        slide, tagline,
        0.3, 2.9, 9.0, 0.8,
        font_name=theme.font_body,
        font_size=18, italic=True,
        color=theme.accent_gold, align="left"
    )

    # Separator line
    add_accent_bar(slide, 0.3, 3.75, 6.0, 0.02, theme.accent_blue)

    # Studio + metadata
    meta = f"{studio_name}  ·  {genre}  ·  {platform}"
    add_text_box(
        slide, meta,
        0.3, 3.9, 9.0, 0.5,
        font_name=theme.font_body,
        font_size=12, color=theme.text_secondary, align="left"
    )

    # "Game Design Document" label in corner
    add_text_box(
        slide, "GAME DESIGN DOCUMENT",
        7.0, 5.1, 2.8, 0.4,
        font_name=theme.font_body,
        font_size=8, color=theme.text_secondary, align="right"
    )


def build_content_slide(
    prs: "Presentation",
    slide_number: int,
    title: str,
    content_blocks: List[Dict[str, Any]],
    theme: PitchTheme = DEFAULT_THEME
) -> None:
    """
    Build a standard content slide with title and content blocks.

    content_blocks: list of dicts with keys:
        - type: "bullets" | "text" | "stat"
        - title: optional sub-heading
        - items: list of strings (for bullets)
        - text: string (for text type)
        - value: string (for stat type — large number display)
        - label: string (for stat type — label below value)
    """
    if not PPTX_AVAILABLE:
        return
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_background(slide, theme.bg_slide)

    # Top accent bar
    add_accent_bar(slide, 0, 0, 10.0, 0.08, theme.accent_blue)

    # Slide number indicator
    add_text_box(
        slide, str(slide_number),
        9.5, 0.1, 0.4, 0.35,
        font_size=8, color=theme.text_secondary, align="right"
    )

    # Title
    add_text_box(
        slide, title,
        0.4, 0.15, 9.0, 0.7,
        font_name=theme.font_heading,
        font_size=22, bold=True,
        color=theme.text_primary, align="left"
    )

    # Divider
    add_accent_bar(slide, 0.4, 0.88, 9.2, 0.02, theme.accent_blue)

    # Content blocks — auto-layout up to 4 blocks in a grid
    n = len(content_blocks)
    if n == 0:
        return
    elif n <= 2:
        col_width = 4.4
        positions = [(0.4, 1.0), (5.2, 1.0)]
    elif n <= 3:
        col_width = 2.9
        positions = [(0.4, 1.0), (3.55, 1.0), (6.7, 1.0)]
    else:
        col_width = 2.1
        positions = [(0.4, 1.0), (2.65, 1.0), (4.9, 1.0), (7.15, 1.0)]

    for i, block in enumerate(content_blocks[:4]):
        if i >= len(positions):
            break
        x, y = positions[i]

        block_type = block.get("type", "bullets")

        if block_type == "stat":
            # Large stat display
            add_text_box(
                slide, block.get("value", "—"),
                x, y, col_width, 1.5,
                font_size=36, bold=True,
                color=theme.accent_gold, align="center"
            )
            add_text_box(
                slide, block.get("label", ""),
                x, y + 1.5, col_width, 0.6,
                font_size=10, color=theme.text_secondary, align="center"
            )

        elif block_type == "text":
            sub_title = block.get("title", "")
            text = block.get("text", "")
            if sub_title:
                add_text_box(
                    slide, sub_title,
                    x, y, col_width, 0.45,
                    font_size=13, bold=True,
                    color=theme.accent_gold, align="left"
                )
                y += 0.45
            add_text_box(
                slide, text,
                x, y, col_width, 3.5,
                font_size=10.5,
                color=theme.text_primary, align="left"
            )

        elif block_type == "bullets":
            sub_title = block.get("title", "")
            items = block.get("items", [])
            add_bullet_text_box(
                slide, sub_title, items,
                x, y, col_width, 4.0,
                theme=theme,
                title_size=13, bullet_size=10
            )


def build_comparison_slide(
    prs: "Presentation",
    slide_number: int,
    title: str,
    competitors: List[Dict[str, str]],
    our_game_name: str,
    theme: PitchTheme = DEFAULT_THEME,
    differentiator: str = ""
) -> None:
    """
    Build a competitive comparison slide.

    competitors: list of dicts with keys:
        - name: competitor title
        - strength: what they do well
        - weakness: what they lack
    differentiator: text explaining how our game is different (displayed at bottom)
    """
    if not PPTX_AVAILABLE:
        return
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_background(slide, theme.bg_slide)

    add_accent_bar(slide, 0, 0, 10.0, 0.08, theme.accent_blue)

    add_text_box(
        slide, str(slide_number),
        9.5, 0.1, 0.4, 0.35,
        font_size=8, color=theme.text_secondary, align="right"
    )
    add_text_box(
        slide, title,
        0.4, 0.15, 9.0, 0.7,
        font_name=theme.font_heading,
        font_size=22, bold=True,
        color=theme.text_primary, align="left"
    )
    add_accent_bar(slide, 0.4, 0.88, 9.2, 0.02, theme.accent_blue)

    # Competitor cards
    n_comps = min(len(competitors), 3)
    comp_width = 2.8
    positions_x = [0.4, 3.5, 6.6]

    for i, comp in enumerate(competitors[:n_comps]):
        x = positions_x[i]
        # Card background
        add_accent_bar(slide, x, 1.0, comp_width, 3.5, theme.bg_card)

        # Competitor name
        add_text_box(
            slide, comp.get("name", ""),
            x + 0.1, 1.05, comp_width - 0.2, 0.5,
            font_size=13, bold=True,
            color=theme.accent_blue, align="center"
        )
        # Strength
        add_text_box(
            slide, f"✓ {comp.get('strength', '')}",
            x + 0.1, 1.65, comp_width - 0.2, 1.0,
            font_size=9.5,
            color=theme.accent_green, align="left"
        )
        # Weakness
        add_text_box(
            slide, f"✗ {comp.get('weakness', '')}",
            x + 0.1, 2.75, comp_width - 0.2, 1.0,
            font_size=9.5,
            color=theme.accent_red, align="left"
        )

    # "Our game" differentiation at bottom
    add_accent_bar(slide, 0.4, 4.7, 9.2, 0.02, theme.accent_gold)
    diff_text = differentiator if differentiator else "[differentiator goes here]"
    add_text_box(
        slide, f"★  {our_game_name}: {diff_text}",
        0.4, 4.85, 9.2, 0.6,
        font_size=12, bold=True,
        color=theme.accent_gold, align="center"
    )


def build_closing_slide(
    prs: "Presentation",
    slide_number: int,
    game_title: str,
    ask: str,
    contact_info: str,
    theme: PitchTheme = DEFAULT_THEME
) -> None:
    """Build the closing/CTA slide."""
    if not PPTX_AVAILABLE:
        return
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_background(slide, theme.bg_dark)

    add_accent_bar(slide, 0, 0, 10.0, 0.08, theme.accent_gold)

    add_text_box(
        slide, game_title,
        0.5, 1.0, 9.0, 1.2,
        font_name=theme.font_heading,
        font_size=36, bold=True,
        color=theme.text_primary, align="center"
    )

    add_text_box(
        slide, ask,
        0.5, 2.2, 9.0, 1.2,
        font_size=16,
        color=theme.accent_gold, align="center"
    )

    add_accent_bar(slide, 2.5, 3.5, 5.0, 0.02, theme.accent_blue)

    add_text_box(
        slide, contact_info,
        0.5, 3.7, 9.0, 0.8,
        font_size=11,
        color=theme.text_secondary, align="center"
    )

    add_text_box(
        slide, str(slide_number),
        9.5, 0.1, 0.4, 0.35,
        font_size=8, color=theme.text_secondary, align="right"
    )

    add_accent_bar(slide, 0, theme.slide_height - 0.08, 10.0, 0.08, theme.accent_gold)
