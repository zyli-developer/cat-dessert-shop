"""
generate_pitch_deck_pptx.py
----------------------------
Generates a professional 10-12 slide pitch deck PowerPoint (.pptx) for a game concept.

Usage:
    python scripts/generate_pitch_deck_pptx.py --title "My Game" --output "MyGame_Pitch.pptx"
    python scripts/generate_pitch_deck_pptx.py --config pitch_content.json --output "MyGame_Pitch.pptx"

Requirements:
    pip install python-pptx

Slide Structure:
  1.  Title Slide (game name, tagline, studio, genre/platform)
  2.  The Problem / Market Opportunity
  3.  The Solution / Game Concept
  4.  Core Gameplay Loop
  5.  Key Features (4 pillars)
  6.  Target Audience & Market Size
  7.  Monetization Strategy
  8.  Competitive Landscape
  9.  Development Timeline
  10. The Team
  11. The Ask / Next Steps
"""

import argparse
import json
import os
import sys
from datetime import datetime
from typing import Dict, List, Any, Optional

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, SCRIPT_DIR)

try:
    from utils.pptx_builder import (
        PitchTheme, DEFAULT_THEME,
        create_presentation,
        add_background, add_text_box, add_bullet_text_box,
        add_accent_bar, add_divider_line,
        build_title_slide, build_content_slide,
        build_comparison_slide, build_closing_slide
    )
    PPTX_BUILDER_AVAILABLE = True
except ImportError:
    PPTX_BUILDER_AVAILABLE = False

try:
    from utils.section_registry import validate_data_sensibility
    REGISTRY_AVAILABLE = True
except ImportError:
    REGISTRY_AVAILABLE = False

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


# ─────────────────────────────────────────────
# DEFAULT CONTENT TEMPLATES
# ─────────────────────────────────────────────

def _get_default_pitch_content(game_data: Dict) -> Dict:
    """Generate default pitch content from game data when specific content isn't provided."""
    title = game_data.get("game_title", "UNTITLED GAME")
    genre = game_data.get("genre", "Genre TBD")
    platform = game_data.get("platform", "PC, Console, Mobile")
    audience = game_data.get("audience", "Ages 18-35, midcore gamers")
    tagline = game_data.get("tagline", f"The definitive {genre} experience")
    monetization = game_data.get("monetization", "Free-to-Play with Battle Pass")
    studio = game_data.get("studio_name", "Studio Name")
    team_size = game_data.get("team_size", "6-10 developers")
    comparable_1 = game_data.get("comparable_1", "Game A")
    comparable_2 = game_data.get("comparable_2", "Game B")
    unique_hook = game_data.get("unique_hook", "[Unique differentiator — complete in GDD interview]")

    return {
        "slide_2_problem": {
            "title": "The Market Opportunity",
            "blocks": [
                {
                    "type": "stat",
                    "value": "$X.XB",
                    "label": f"SOURCE NEEDED: {genre} market size (year)"
                },
                {
                    "type": "stat",
                    "value": "XXM",
                    "label": "SOURCE NEEDED: Active players in genre (Steam + mobile)"
                },
                {
                    "type": "bullets",
                    "title": "The Gap",
                    "items": [
                        f"Players want {genre} but with [missing element]",
                        "Existing titles haven't [unmet need]",
                        "Last breakout title in genre: [Title, Year] — gap since then",
                        "Our target audience actively requests [feature] in forums"
                    ]
                }
            ]
        },
        "slide_3_solution": {
            "title": f"The Game: {title}",
            "blocks": [
                {
                    "type": "text",
                    "title": "High Concept",
                    "text": tagline
                },
                {
                    "type": "text",
                    "title": "Core Fantasy",
                    "text": f"Players feel like [power fantasy] — the experience {comparable_1} promised but never delivered."
                },
                {
                    "type": "bullets",
                    "title": "Three Pillars",
                    "items": [
                        "[Pillar 1: core mechanic experience]",
                        "[Pillar 2: progression/mastery]",
                        "[Pillar 3: social/session hook]"
                    ]
                }
            ]
        },
        "slide_4_loop": {
            "title": "Core Gameplay Loop",
            "blocks": [
                {
                    "type": "text",
                    "title": "Micro Loop (2–5 min)",
                    "text": "[Describe what the player does in a single combat/puzzle/interaction cycle]"
                },
                {
                    "type": "text",
                    "title": "Macro Loop (20–40 min)",
                    "text": "[Describe how micro loops build into a session — what the player achieves in one sitting]"
                },
                {
                    "type": "text",
                    "title": "Meta Loop (Days → Weeks)",
                    "text": "[Describe the long-term progression arc that creates return behavior]"
                }
            ]
        },
        "slide_5_features": {
            "title": "Key Features",
            "blocks": [
                {
                    "type": "bullets",
                    "title": "★ [Feature 1 Name]",
                    "items": ["[One sentence describing the specific mechanic or system]"]
                },
                {
                    "type": "bullets",
                    "title": "★ [Feature 2 Name]",
                    "items": ["[One sentence — the innovation that doesn't exist elsewhere]"]
                },
                {
                    "type": "bullets",
                    "title": "★ [Feature 3 Name]",
                    "items": ["[One sentence describing player-facing hook]"]
                },
                {
                    "type": "bullets",
                    "title": "★ [Feature 4 Name]",
                    "items": ["[One sentence — the social or retention hook]"]
                }
            ]
        },
        "slide_6_audience": {
            "title": "Target Audience & Market",
            "blocks": [
                {
                    "type": "bullets",
                    "title": "Primary Audience",
                    "items": [
                        audience,
                        "Platform: " + platform,
                        "[Hrs/week in similar games: estimate]",
                        "[Income range / spending behavior]"
                    ]
                },
                {
                    "type": "bullets",
                    "title": "Comparable Player Bases",
                    "items": [
                        f"{comparable_1}: [X]M players, [MAU data]",
                        f"{comparable_2}: [X]M players, [MAU data]",
                        "We target [%] crossover from each"
                    ]
                },
                {
                    "type": "stat",
                    "value": "D30",
                    "label": "SOURCE NEEDED: Retention Target: >X% (cite industry benchmark)"
                }
            ]
        },
        "slide_7_monetization": {
            "title": "Monetization Strategy",
            "blocks": [
                {
                    "type": "bullets",
                    "title": "Revenue Model",
                    "items": [
                        f"Model: {monetization}",
                        "Battle Pass: $9.99/season (8 weeks)",
                        "Cosmetics IAP: $0.99 – $9.99",
                        "No pay-to-win — cosmetics only"
                    ]
                },
                {
                    "type": "stat",
                    "value": "$XX",
                    "label": "SOURCE NEEDED: Target blended ARPU (monthly)"
                },
                {
                    "type": "stat",
                    "value": "X–X%",
                    "label": "SOURCE NEEDED: D30 conversion rate target"
                }
            ]
        },
        "slide_8_competitive": {
            "title": "Competitive Landscape",
            "competitors": [
                {
                    "name": comparable_1,
                    "strength": "[What they do well — 1 sentence]",
                    "weakness": "[What they lack that we deliver — 1 sentence]"
                },
                {
                    "name": comparable_2,
                    "strength": "[Their strength]",
                    "weakness": "[Their gap we fill]"
                },
                {
                    "name": "[Indirect Comp]",
                    "strength": "[Why players go there]",
                    "weakness": "[Why they'd prefer us]"
                }
            ]
        },
        "slide_9_timeline": {
            "title": "Development Timeline",
            "blocks": [
                {
                    "type": "bullets",
                    "title": "Milestones",
                    "items": [
                        "Q[X] 202X — Vertical Slice (core loop playable)",
                        "Q[X] 202X — Alpha (feature complete, internal test)",
                        "Q[X] 202X — Beta (content complete, external playtest)",
                        "Q[X] 202X — Soft Launch (limited regions)",
                        "Q[X] 202X — Global Launch"
                    ]
                },
                {
                    "type": "bullets",
                    "title": "Current Status",
                    "items": [
                        "[Current milestone and % complete]",
                        "[Team size and key roles filled]",
                        "[Budget status: self-funded / seeking investment]",
                        "[Prototype playable: yes/no]"
                    ]
                }
            ]
        },
        "slide_10_team": {
            "title": "The Team",
            "blocks": [
                {
                    "type": "bullets",
                    "title": "Core Team",
                    "items": [
                        f"[Name] — Game Director / Lead Designer: [Credential — shipped title or years exp]",
                        "[Name] — Technical Lead: [Credential]",
                        "[Name] — Art Director: [Credential]",
                        "[Name] — Producer: [Credential]"
                    ]
                },
                {
                    "type": "bullets",
                    "title": "Team Profile",
                    "items": [
                        f"Total: {team_size}",
                        "[Combined years exp / shipped titles]",
                        "[Notable studio backgrounds, if any]",
                        "[Key advisors, if any]"
                    ]
                }
            ]
        },
        "slide_11_ask": {
            "ask": "[The Ask: Seeking $X for Y months runway to Z milestone]",
            "contact": f"{studio}  ·  contact@studio.com  ·  studio.gg"
        }
    }


# ─────────────────────────────────────────────
# MAIN GENERATION FUNCTION
# ─────────────────────────────────────────────

def generate_pitch_deck(
    game_data: Dict,
    output_path: str,
    theme: Optional["PitchTheme"] = None,
    strict: bool = False
) -> str:
    """
    Generate a complete pitch deck .pptx file.

    Args:
        game_data: Dictionary with game metadata and optional slide content.
        output_path: Output .pptx file path.
        theme: Optional PitchTheme. Defaults to DEFAULT_THEME.
        strict: If True, fail export if unsourced metrics or placeholders remain.

    Returns:
        Absolute path to the generated file.
    """
    if not PPTX_AVAILABLE:
        raise ImportError(
            "python-pptx is required. Install with:\n"
            "  pip install python-pptx"
        )
    if not PPTX_BUILDER_AVAILABLE:
        raise ImportError("pptx_builder utils not found. Check scripts/utils/pptx_builder.py")

    # Pre-export validation: check sections if provided
    if REGISTRY_AVAILABLE:
        sections_to_validate = game_data.get("sections", {})
        if sections_to_validate:
            sensibility_warnings = validate_data_sensibility(
                sections_to_validate, strict=strict
            )
            for warning in sensibility_warnings:
                print(f"  WARNING: {warning}")
            if strict and sensibility_warnings:
                raise SystemExit(
                    "STRICT MODE: Export aborted due to unsourced metrics or "
                    "placeholders in business sections. Fix the warnings above "
                    "or remove --strict to export with warnings."
                )

    # Check pitch slide content for SOURCE NEEDED placeholders
    import re
    source_needed_re = re.compile(r"\bSOURCE NEEDED\b", re.IGNORECASE)
    pitch_slides = game_data.get("pitch_slides", {})
    pitch_warnings = []
    for slide_key, slide_data in pitch_slides.items():
        slide_text = json.dumps(slide_data)
        if source_needed_re.search(slide_text):
            pitch_warnings.append(
                f"PLACEHOLDER in '{slide_key}': Contains 'SOURCE NEEDED' markers. "
                f"Replace with real data before external use."
            )
    for warning in pitch_warnings:
        print(f"  WARNING: {warning}")
    if strict and pitch_warnings:
        raise SystemExit(
            "STRICT MODE: Export aborted due to SOURCE NEEDED placeholders "
            "in pitch slides. Fix the warnings above or remove --strict."
        )

    if theme is None:
        theme = DEFAULT_THEME

    prs = create_presentation(theme)

    # Merge default content with provided content
    default_content = _get_default_pitch_content(game_data)
    pitch_content = {**default_content, **game_data.get("pitch_slides", {})}

    game_title = game_data.get("game_title", "UNTITLED GAME")
    tagline = game_data.get("tagline", "A new gaming experience")
    studio = game_data.get("studio_name", "Studio Name")
    genre = game_data.get("genre", "Genre")
    platform = game_data.get("platform", "Platform")

    slide_num = 1

    # ── Slide 1: Title ──────────────────────────────────────
    build_title_slide(
        prs, game_title, tagline, studio, genre, platform, theme
    )
    slide_num += 1

    # ── Slide 2: Market Opportunity ──────────────────────────
    s2 = pitch_content.get("slide_2_problem", {})
    build_content_slide(
        prs, slide_num,
        s2.get("title", "The Market Opportunity"),
        s2.get("blocks", []),
        theme
    )
    slide_num += 1

    # ── Slide 3: Game Concept ────────────────────────────────
    s3 = pitch_content.get("slide_3_solution", {})
    build_content_slide(
        prs, slide_num,
        s3.get("title", f"The Game: {game_title}"),
        s3.get("blocks", []),
        theme
    )
    slide_num += 1

    # ── Slide 4: Core Loop ───────────────────────────────────
    s4 = pitch_content.get("slide_4_loop", {})
    build_content_slide(
        prs, slide_num,
        s4.get("title", "Core Gameplay Loop"),
        s4.get("blocks", []),
        theme
    )
    slide_num += 1

    # ── Slide 5: Key Features ────────────────────────────────
    s5 = pitch_content.get("slide_5_features", {})
    build_content_slide(
        prs, slide_num,
        s5.get("title", "Key Features"),
        s5.get("blocks", []),
        theme
    )
    slide_num += 1

    # ── Slide 6: Audience & Market ───────────────────────────
    s6 = pitch_content.get("slide_6_audience", {})
    build_content_slide(
        prs, slide_num,
        s6.get("title", "Target Audience & Market"),
        s6.get("blocks", []),
        theme
    )
    slide_num += 1

    # ── Slide 7: Monetization ────────────────────────────────
    s7 = pitch_content.get("slide_7_monetization", {})
    build_content_slide(
        prs, slide_num,
        s7.get("title", "Monetization Strategy"),
        s7.get("blocks", []),
        theme
    )
    slide_num += 1

    # ── Slide 8: Competitive Landscape ───────────────────────
    s8 = pitch_content.get("slide_8_competitive", {})
    competitors = s8.get("competitors", [
        {"name": "Comp 1", "strength": "Strong brand", "weakness": "No mobile version"},
        {"name": "Comp 2", "strength": "Large player base", "weakness": "Pay-to-win criticism"},
        {"name": "Comp 3", "strength": "Deep mechanics", "weakness": "High skill barrier"},
    ])
    differentiator = s8.get("differentiator", game_data.get("unique_hook", ""))
    build_comparison_slide(
        prs, slide_num,
        s8.get("title", "Competitive Landscape"),
        competitors,
        game_title,
        theme,
        differentiator=differentiator
    )
    slide_num += 1

    # ── Slide 9: Development Timeline ────────────────────────
    s9 = pitch_content.get("slide_9_timeline", {})
    build_content_slide(
        prs, slide_num,
        s9.get("title", "Development Timeline"),
        s9.get("blocks", []),
        theme
    )
    slide_num += 1

    # ── Slide 10: Team ────────────────────────────────────────
    s10 = pitch_content.get("slide_10_team", {})
    build_content_slide(
        prs, slide_num,
        s10.get("title", "The Team"),
        s10.get("blocks", []),
        theme
    )
    slide_num += 1

    # ── Slide 11: The Ask / Closing ───────────────────────────
    s11 = pitch_content.get("slide_11_ask", {})
    # Support both simple ask/contact format and structured blocks format
    ask_text = s11.get("ask", "")
    if not ask_text and "blocks" in s11:
        # Extract ask text from blocks: find stat blocks and bullet blocks
        parts = []
        for block in s11["blocks"]:
            if block.get("type") == "stat":
                parts.append(f"{block.get('value', '')} — {block.get('label', '')}")
            elif block.get("type") == "bullets":
                block_title = block.get("title", "")
                items = block.get("items", [])
                if block_title:
                    parts.append(block_title + ": " + "; ".join(items[:3]))
        ask_text = "\n".join(parts) if parts else "[The Ask — investment amount and milestone target]"
    elif not ask_text:
        ask_text = "[The Ask — investment amount and milestone target]"

    build_closing_slide(
        prs, slide_num,
        game_title,
        ask_text,
        s11.get("contact", f"{studio}  ·  contact@studio.com"),
        theme
    )

    # Save
    output_dir = os.path.dirname(os.path.abspath(output_path))
    os.makedirs(output_dir, exist_ok=True)
    prs.save(output_path)
    abs_path = os.path.abspath(output_path)
    print(f"✓ Pitch deck generated: {abs_path}")
    print(f"  Slides: {slide_num}")
    return abs_path


# ─────────────────────────────────────────────
# CLI ENTRY POINT
# ─────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(
        description="Generate a game pitch deck (.pptx)",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  python generate_pitch_deck_pptx.py --title "Echo Chamber" --output "EchoChamber_Pitch.pptx"
  python generate_pitch_deck_pptx.py --config pitch_content.json --output "MyGame_Pitch.pptx"
        """
    )
    parser.add_argument("--title", help="Game title")
    parser.add_argument("--studio", default="Studio Name", help="Studio name")
    parser.add_argument("--genre", default="Genre TBD", help="Genre")
    parser.add_argument("--platform", default="PC, Console, Mobile", help="Platform")
    parser.add_argument("--audience", default="Ages 18-35, midcore gamers", help="Audience")
    parser.add_argument("--tagline", default="", help="Tagline")
    parser.add_argument("--config", help="JSON config file with pitch content")
    parser.add_argument("--output", default="pitch_deck.pptx", help="Output .pptx path")
    parser.add_argument("--strict", action="store_true",
                        help="Fail export if SOURCE NEEDED placeholders or unsourced metrics remain")

    args = parser.parse_args()

    if not PPTX_AVAILABLE:
        print("ERROR: python-pptx is not installed. Install with: pip install python-pptx")
        sys.exit(1)

    if args.config:
        try:
            with open(args.config, "r", encoding="utf-8") as f:
                game_data = json.load(f)
        except (FileNotFoundError, json.JSONDecodeError) as e:
            print(f"ERROR loading config: {e}")
            sys.exit(1)
    elif args.title:
        game_data = {
            "game_title": args.title,
            "tagline": args.tagline or f"A new {args.genre} experience",
            "genre": args.genre,
            "platform": args.platform,
            "audience": args.audience,
            "studio_name": args.studio,
        }
    else:
        parser.print_help()
        print("\nERROR: Provide --title or --config")
        sys.exit(1)

    try:
        generate_pitch_deck(game_data=game_data, output_path=args.output, strict=args.strict)
    except Exception as e:
        print(f"ERROR: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)


if __name__ == "__main__":
    main()
